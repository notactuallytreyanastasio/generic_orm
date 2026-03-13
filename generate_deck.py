#!/usr/bin/env python3
"""Generate deck.pptx matching the reveal.js HTML deck styling."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK = RGBColor(0x11, 0x11, 0x11)
BODY = RGBColor(0x55, 0x55, 0x55)
LABEL = RGBColor(0x99, 0x99, 0x99)
LIGHT = RGBColor(0xAA, 0xAA, 0xAA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xDD, 0xDD, 0xDD)
LIGHTBG = RGBColor(0xF7, 0xF7, 0xF7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
MARGIN = Inches(0.8)
CONTENT_W = SLIDE_W - 2 * MARGIN


def add_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    return txBox.text_frame


def set_run(para, text, size=18, color=BODY, bold=False, italic=False, font_name="Helvetica Neue"):
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    return run


def add_para(tf, text, size=18, color=BODY, bold=False, italic=False, align=PP_ALIGN.LEFT, space_before=0, space_after=0, font_name="Helvetica Neue"):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    set_run(p, text, size, color, bold, italic, font_name)
    return p


def add_section_header(slide, text):
    tf = add_textbox(slide, MARGIN, Inches(0.5), CONTENT_W, Inches(0.5))
    add_para(tf, text.upper(), size=13, color=LABEL, bold=False, align=PP_ALIGN.LEFT)


def add_title(slide, text, top=Inches(0.9)):
    tf = add_textbox(slide, MARGIN, top, CONTENT_W, Inches(0.8))
    add_para(tf, text, size=36, color=DARK, bold=True, align=PP_ALIGN.LEFT)


def add_body(slide, text, top=Inches(1.8), size=18, left=None, width=None):
    l = left if left else MARGIN
    w = width if width else CONTENT_W
    tf = add_textbox(slide, l, top, w, Inches(1.0))
    add_para(tf, text, size=size, color=BODY)
    return tf


def add_big_num(slide, number, label, left, top, width=Inches(3.0)):
    tf = add_textbox(slide, left, top, width, Inches(1.4))
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_para(tf, str(number), size=60, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
    add_para(tf, label, size=12, color=LABEL, align=PP_ALIGN.CENTER, space_before=4)


def add_table(slide, rows, col_widths, top, left=None):
    l = left if left else MARGIN
    n_rows = len(rows)
    n_cols = len(rows[0])
    total_w = sum(col_widths)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, l, top, Emu(total_w), Inches(n_rows * 0.38))
    tbl = tbl_shape.table

    for ci, w in enumerate(col_widths):
        tbl.columns[ci].width = Emu(w)

    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT

            is_header = ri == 0
            run = p.add_run()
            run.text = str(cell_text)
            run.font.size = Pt(11 if is_header else 13)
            run.font.name = "Helvetica Neue"
            run.font.bold = is_header
            run.font.color.rgb = LABEL if is_header else BODY

            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

    return tbl_shape


def add_rect_box(slide, left, top, width, height, text, text_size=12, text_color=BODY, border_color=BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = border_color
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(text_size)
    run.font.color.rgb = text_color
    run.font.name = "Helvetica Neue"
    return shape


# ── SLIDE 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
tf = add_textbox(slide, MARGIN, Inches(1.8), CONTENT_W, Inches(1.2))
add_para(tf, "Temper", size=72, color=BLACK, bold=True, align=PP_ALIGN.CENTER)

tf2 = add_textbox(slide, MARGIN, Inches(3.2), CONTENT_W, Inches(0.6))
add_para(tf2, "Write Once, Secure Everywhere", size=28, color=LABEL, align=PP_ALIGN.CENTER)

tf3 = add_textbox(slide, MARGIN, Inches(4.0), CONTENT_W, Inches(0.5))
add_para(tf3, "A cross-compiled language for security-critical infrastructure at scale", size=16, color=LIGHT, align=PP_ALIGN.CENTER)

langs = ["JavaScript", "Python", "Rust", "Java", "Lua", "C#", "and more to come..."]
pill_w = Inches(1.45)
pill_h = Inches(0.38)
total_pills_w = len(langs) * pill_w + (len(langs) - 1) * Inches(0.15)
start_x = (SLIDE_W - total_pills_w) // 2
for i, lang in enumerate(langs):
    x = start_x + i * (pill_w + Inches(0.15))
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.2), pill_w, pill_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.75)
    if lang == "and more to come...":
        shape.line.dash_style = 3  # dash
    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = lang
    r.font.size = Pt(11)
    r.font.color.rgb = LABEL
    r.font.name = "Helvetica Neue"

# ── SLIDE 2: The Problem ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "The Problem")
add_title(slide, "SDK maintenance at scale is a security crisis")

gap = Inches(3.6)
start = (SLIDE_W - 3 * gap) // 2
for i, (num, lbl) in enumerate([("3,600+", "AWS service client\nlibraries"), ("12", "language\ntargets"), ("156K+", "operation\nimplementations")]):
    add_big_num(slide, num, lbl, start + i * gap, Inches(2.0), Inches(3.0))

tf = add_body(slide, "", top=Inches(4.0))
p = tf.paragraphs[0]
set_run(p, "A single vulnerability in serialization logic must be patched ", 16, BODY)
set_run(p, "independently in every language", 16, BLACK, bold=True)
set_run(p, ". The Java SDK v1\u2192v2 migration spanned ", 16, BODY)
set_run(p, "7 years", 16, BLACK, bold=True)
set_run(p, ". Three major SDKs hit end-of-support in 2025 \u2014 ", 16, BODY)
set_run(p, "zero security patches after cutoff", 16, BLACK, bold=True)
set_run(p, ".", 16, BODY)

tf2 = add_body(slide, "", top=Inches(4.8))
p2 = tf2.paragraphs[0]
set_run(p2, ".NET SDK v3\u2192v4 has a ", 16, BODY)
set_run(p2, "3-month maintenance window", 16, BLACK, bold=True)
set_run(p2, " vs the 12-month default. The trend is accelerating.", 16, BODY)

tf3 = add_textbox(slide, MARGIN, Inches(6.6), CONTENT_W, Inches(0.4))
add_para(tf3, "Source: AWS SDKs and Tools maintenance policy, March 2026", size=10, color=LIGHT)

# ── SLIDE 3: The Temper Answer ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "The Answer")
add_title(slide, "One source. Six backends. Zero drift.")
add_body(slide, "Temper compiles a single source tree to idiomatic output in 6 languages. Not transpilation \u2014 native, typed, importable libraries.", top=Inches(1.8), size=16)

flow_items = ["Temper Source", "\u2192", "temper build", "\u2192", "ES Modules", "Python 3", "Cargo Crate", "Java", "Lua", ".NET"]
box_w = Inches(1.2)
box_h = Inches(0.35)
flow_y = Inches(2.7)
x = MARGIN
for item in flow_items:
    if item == "\u2192":
        tf = add_textbox(slide, x, flow_y, Inches(0.3), box_h)
        add_para(tf, "\u2192", size=16, color=LIGHT, align=PP_ALIGN.CENTER)
        x += Inches(0.3)
    else:
        add_rect_box(slide, x, flow_y, box_w, box_h, item, text_size=10, border_color=BORDER if item != "Temper Source" else LABEL)
        x += box_w + Inches(0.08)

# Smithy vs Temper cards
card_w = Inches(5.6)
card_h = Inches(2.8)
card_y = Inches(3.5)
card_x1 = MARGIN
card_x2 = MARGIN + card_w + Inches(0.6)

for cx, title, items in [
    (card_x1, "What Smithy does for AWS", [
        "Generates API clients from models",
        "Still requires per-language runtime teams",
        "Security fixes need per-SDK patching",
        "Multi-year major version migrations",
    ]),
    (card_x2, "What Temper does better", [
        "Compiles arbitrary logic, not just API shapes",
        "One codebase = one security surface",
        "Fix once, rebuild, done everywhere",
        "Same-day propagation to all backends",
    ]),
]:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, card_y, card_w, card_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.75)

    tf = add_textbox(slide, cx + Inches(0.25), card_y + Inches(0.2), card_w - Inches(0.5), Inches(0.4))
    add_para(tf, title, size=16, color=DARK, bold=True)

    tf2 = add_textbox(slide, cx + Inches(0.25), card_y + Inches(0.7), card_w - Inches(0.5), Inches(2.0))
    for item in items:
        add_para(tf2, "\u2014  " + item, size=13, color=BODY, space_after=4)

# ── SLIDE 4: Proof - The ORM ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "Proof")
add_title(slide, "A type-safe SQL ORM, running in 6 languages")
add_body(slide, "We built a security-focused ORM with 5 defense layers against SQL injection, compiled it to 6 backends, and built 6 identical todo-list apps to prove it works.", top=Inches(1.8), size=16)

cw = Inches(11.733)
rows = [
    ["Language", "Framework", "ORM Import", "SQL Injection Vulns"],
    ["JavaScript", "Express + EJS", "from 'orm/src'", "0"],
    ["Python", "Flask + Jinja2", "from orm.src import", "0"],
    ["Rust", "Axum + askama", "use orm::src::{...}", "0"],
    ["Java", "Spring Boot", "import orm.src.*", "0"],
    ["Lua", "Raw socket HTTP", 'require("orm/src")', "0"],
    ["C#", "ASP.NET Razor", "using Orm.Src", "0"],
]
col_ws = [Emu(int(cw * r)) for r in [0.2, 0.22, 0.38, 0.2]]
add_table(slide, rows, col_ws, Inches(2.6))

tf = add_textbox(slide, MARGIN, Inches(6.6), CONTENT_W, Inches(0.4))
add_para(tf, "Full MITRE ATT&CK / CWE analysis conducted per app. Reports published in each app repository.", size=10, color=LIGHT)

# ── SLIDE 5: Defense Layers ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "Security Architecture")
add_title(slide, "5 compile-time defense layers")

rows = [
    ["#", "Layer", "CWE", "What It Prevents"],
    ["1", "SafeIdentifier", "CWE-89", "Table/column name injection. Regex-validated, sealed class not exported."],
    ["2", "Typed SqlPart", "CWE-89", "Value injection. Sealed interface with exhaustive type dispatch."],
    ["3", "SqlBuilder", "CWE-89", "Structural injection. No appendRaw. Structure and data strictly separated."],
    ["4", "Changeset", "CWE-915", "Mass assignment. cast() requires List<SafeIdentifier>."],
    ["5", "Query Builder", "CWE-400", "Unbounded queries. WHERE accepts only SqlFragment."],
]
cw = Inches(11.733)
col_ws = [Emu(int(cw * r)) for r in [0.05, 0.18, 0.12, 0.65]]
add_table(slide, rows, col_ws, Inches(2.0))

tf = add_body(slide, "These defenses are structural, not advisory. The type system makes it impossible to construct unsafe SQL without deliberately circumventing the API.", top=Inches(5.6), size=16)

# ── SLIDE 6: Fix Once ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "Fix Once, Fix Everywhere")
add_title(slide, "3 findings. 1 commit. 6 backends patched.")
add_body(slide, "MITRE analysis identified 3 ORM-level concerns. We fixed all 3 in one Temper commit and rebuilt.", top=Inches(1.8), size=16)

rows = [
    ["Finding", "Severity", "Issue", "Fix", "Status"],
    ["ORM-1", "Medium", "Column names passed as String", "Routed through SafeIdentifier", "Resolved"],
    ["ORM-2", "Low", "SqlDate missing quote escaping", "Char-by-char escaping", "Resolved"],
    ["ORM-3", "Low", "SqlFloat64 produces NaN/Infinity", "Renders NULL", "Resolved"],
    ["ORM-4", "Info", "Escaping-based, not parameterized", "Design limitation", "Acknowledged"],
]
cw = Inches(11.733)
col_ws = [Emu(int(cw * r)) for r in [0.1, 0.1, 0.32, 0.3, 0.18]]
add_table(slide, rows, col_ws, Inches(2.5))

# Code block
code_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(4.8), CONTENT_W, Inches(1.2))
code_shape.fill.solid()
code_shape.fill.fore_color.rgb = LIGHTBG
code_shape.line.color.rgb = BORDER
code_shape.line.width = Pt(0.5)
ctf = code_shape.text_frame
ctf.word_wrap = True
p = ctf.paragraphs[0]
set_run(p, "// One change in Temper source (changeset.temper.md)", 11, LIGHT, font_name="Courier New")
p2 = ctf.add_paragraph()
set_run(p2, "colNames.add(pair.key);", 11, LIGHT, font_name="Courier New")
set_run(p2, "           // raw String", 11, LIGHT, font_name="Courier New")
p3 = ctf.add_paragraph()
set_run(p3, "colNames.add(fd.name.sqlValue);", 11, BLACK, bold=True, font_name="Courier New")
set_run(p3, "   // SafeIdentifier \u2192 validated", 11, BODY, font_name="Courier New")

tf = add_body(slide, "", top=Inches(6.2))
p = tf.paragraphs[0]
set_run(p, "Traditional multi-SDK: ", 16, BODY)
set_run(p, "6 PRs, 6 reviews, 6 releases", 16, BLACK, bold=True)
set_run(p, ". Temper: ", 16, BODY)
set_run(p, "1 commit, 1 build, done.", 16, BLACK, bold=True)

# ── SLIDE 7: Head-to-Head ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "AI-Generated ORM Comparison")
add_title(slide, "Same prompts. Same goals. Different security.")
add_body(slide, "Two ORMs built with comparable AI prompts (3\u20134 prompts, same goals). One in Ruby, one in Temper. Both underwent identical MITRE/CWE analysis.", top=Inches(1.8), size=16)

# F vs A
vs_y = Inches(2.6)
left_center = SLIDE_W // 4
right_center = SLIDE_W * 3 // 4

for cx, grade, name, lang, subtitle in [
    (left_center - Inches(1.5), "F", "SomeORM", "(Ruby)", "CVSS 9.8"),
    (right_center - Inches(1.5), "A", "Generic ORM", "(Temper)", "0 injection vulnerabilities"),
]:
    tf = add_textbox(slide, cx, vs_y, Inches(3), Inches(0.8))
    add_para(tf, grade, size=48, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
    tf2 = add_textbox(slide, cx, vs_y + Inches(0.8), Inches(3), Inches(0.4))
    p = tf2.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p, name + " ", 18, BLACK, bold=True)
    set_run(p, lang, 18, LABEL)
    tf3 = add_textbox(slide, cx, vs_y + Inches(1.2), Inches(3), Inches(0.3))
    add_para(tf3, subtitle, size=12, color=BODY, align=PP_ALIGN.CENTER)

tf_vs = add_textbox(slide, SLIDE_W // 2 - Inches(0.3), vs_y + Inches(0.2), Inches(0.6), Inches(0.5))
add_para(tf_vs, "vs", size=20, color=BORDER, align=PP_ALIGN.CENTER)

rows = [
    ["Metric", "SomeORM (Ruby)", "Generic ORM (Temper)"],
    ["SQL injection vulnerabilities", "9 critical/high", "0"],
    ["Identifier validation", "None \u2014 direct interpolation", "SafeIdentifier sealed type"],
    ["Value escaping", "Single-quote only", "Typed SqlPart hierarchy (6 types)"],
    ["Mass assignment protection", "None", "Changeset cast() whitelist"],
    ["Languages supported", "1", "6"],
    ["Fix propagation", "Manual", "Automatic \u2014 1 commit"],
]
cw = Inches(11.733)
col_ws = [Emu(int(cw * r)) for r in [0.3, 0.35, 0.35]]
add_table(slide, rows, col_ws, Inches(4.4))

# ── SLIDE 8: SomeORM 9 Vulns ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "Conventional ORM: 9 Attack Vectors")
add_title(slide, "What AI-generated code looks like without Temper")

rows = [
    ["#", "CWE", "Attack Surface", "CVSS", "Exploit"],
    ["1", "CWE-89", "Table name injection", "9.8", 'from("users; DROP TABLE users; --")'],
    ["2", "CWE-89", "SELECT field injection", "9.1", 'select("(SELECT password FROM admin)")'],
    ["3", "CWE-89", "WHERE raw SQL", "9.8", 'where("id = #{user_input}")'],
    ["4", "CWE-89", "Comparison field injection", "8.6", 'eq("id) OR 1=1; --", 1)'],
    ["5", "CWE-89", "GROUP BY injection", "8.1", 'group_by("user_id; DROP TABLE")'],
    ["6", "CWE-89", "ORDER BY injection", "8.1", 'order_by("id; DROP TABLE")'],
    ["7", "CWE-89", "Direction injection", "5.3", 'order_by("id", direction: "ASC; --")'],
    ["8", "CWE-20", "LIMIT type confusion", "5.3", 'limit("10; DROP TABLE")'],
    ["9", "CWE-116", "Incomplete escaping", "6.5", "Backslash, NULL byte, encoding"],
]
cw = Inches(11.733)
col_ws = [Emu(int(cw * r)) for r in [0.04, 0.09, 0.22, 0.07, 0.58]]
add_table(slide, rows, col_ws, Inches(1.9))

tf = add_body(slide, "", top=Inches(6.3))
p = tf.paragraphs[0]
set_run(p, "Every one of these is ", 16, BODY)
set_run(p, "structurally impossible", 16, BLACK, bold=True)
set_run(p, " in the Temper ORM. Prevented by the type system, not developer discipline.", 16, BODY)

# ── SLIDE 9: AI Development ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "AI-First Development")
add_title(slide, "Temper makes AI-generated code secure by default")
add_body(slide, "The question is not whether AI writes your code. It is whether the language constrains AI to produce safe output.", top=Inches(1.8), size=16)

card_w = Inches(5.6)
card_h = Inches(2.8)
card_y = Inches(2.6)

for cx, title, items in [
    (MARGIN, "Without Temper", [
        "AI generates plausible but insecure SQL builders",
        "String concatenation is the path of least resistance",
        "Escaping bugs are subtle and language-specific",
        "Each language needs independent security review",
        "9 critical vulns from 3\u20134 prompts (SomeORM)",
    ]),
    (MARGIN + card_w + Inches(0.6), "With Temper", [
        "Type system prevents unsafe SQL construction",
        "SafeIdentifier + sealed interfaces = no bypass",
        "AI cannot accidentally skip validation",
        "One security review covers all 6 backends",
        "0 vulnerabilities from same-effort prompts",
    ]),
]:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, card_y, card_w, card_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.75)

    tf = add_textbox(slide, cx + Inches(0.25), card_y + Inches(0.2), card_w - Inches(0.5), Inches(0.4))
    add_para(tf, title, size=16, color=DARK, bold=True)

    tf2 = add_textbox(slide, cx + Inches(0.25), card_y + Inches(0.7), card_w - Inches(0.5), Inches(2.0))
    for item in items:
        add_para(tf2, "\u2014  " + item, size=13, color=BODY, space_after=4)

# Quote
quote_y = Inches(5.8)
# Left border line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, quote_y, Inches(0.04), Inches(0.7))
line.fill.solid()
line.fill.fore_color.rgb = BORDER
line.line.fill.background()

tf = add_textbox(slide, MARGIN + Inches(0.25), quote_y, CONTENT_W - Inches(0.25), Inches(0.7))
add_para(tf, "The best security does not depend on the developer \u2014 or the AI \u2014 remembering to do the right thing. Temper makes the wrong thing impossible to express.", size=14, color=LABEL, italic=True)

# ── SLIDE 9b: Built by AI ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "Case Study")
add_title(slide, "This entire project was built by AI")
add_body(slide, "Everything in this deck was produced by a single developer working with Claude. The ORM, the apps, the security analysis, the fixes, and this deck itself.", top=Inches(1.8), size=16)

rows = [
    ["What", "How", "Result"],
    ["ORM source (Temper)", "AI-generated with iterative prompts", "5 defense layers, 56 tests, 0 vulns"],
    ["6 todo-list applications", "AI-generated per language", "6 frameworks, identical functionality"],
    ["CI/CD pipeline", "AI-generated GitHub Actions", "12 repos, automated cascade"],
    ["MITRE/CWE security audit", "AI-conducted analysis", "4 ORM findings, 0 app-level injection"],
    ["Bug fixes (ORM-1/2/3)", "AI-identified and AI-fixed", "1 commit, 6 backends patched"],
    ["Verification", "AI grepped compiled output", "Fixes confirmed in all backends"],
]
cw = Inches(11.733)
col_ws = [Emu(int(cw * r)) for r in [0.28, 0.38, 0.34]]
add_table(slide, rows, col_ws, Inches(2.6))

tf = add_body(slide, "", top=Inches(6.0))
p = tf.paragraphs[0]
set_run(p, "The AI did not need to know 6 languages to secure 6 backends. It worked in ", 16, BODY)
set_run(p, "one language", 16, BLACK, bold=True)
set_run(p, " and the compiler did the rest. This is the workflow AI-driven development needs.", 16, BODY)

# ── SLIDE 9c: AI Agents + Temper ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "The AI Workflow")
add_title(slide, "Why AI agents need typed, cross-compiled foundations")
add_body(slide, "As AI coding agents move from autocomplete to autonomous engineering, the infrastructure they build on determines whether their output is safe at scale.", top=Inches(1.8), size=16)

card_w = Inches(5.6)
card_h = Inches(3.0)
card_y = Inches(2.6)

for cx, title, items in [
    (MARGIN, "The problem with AI + conventional languages", [
        "AI generates idiomatic code \u2014 including idiomatic mistakes",
        "Ruby, Python, JS all allow raw string SQL with zero friction",
        "An agent building SDKs must learn security patterns per language",
        "Review burden multiplies: N languages \u00d7 M agents \u00d7 K services",
        "A single missed escaping call in one language = a CVE",
    ]),
    (MARGIN + card_w + Inches(0.6), "The advantage of AI + Temper", [
        "AI writes to a single type system with enforced invariants",
        "Sealed interfaces make entire vulnerability classes unrepresentable",
        "Agents work at the abstraction level, not the language level",
        "Review burden is constant: 1 source, regardless of target count",
        "Adding a 7th backend target requires zero security re-audit",
    ]),
]:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, card_y, card_w, card_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.75)

    tf = add_textbox(slide, cx + Inches(0.25), card_y + Inches(0.2), card_w - Inches(0.5), Inches(0.4))
    add_para(tf, title, size=14, color=DARK, bold=True)

    tf2 = add_textbox(slide, cx + Inches(0.25), card_y + Inches(0.7), card_w - Inches(0.5), Inches(2.2))
    for item in items:
        add_para(tf2, "\u2014  " + item, size=12, color=BODY, space_after=4)

tf = add_body(slide, "", top=Inches(6.0))
p = tf.paragraphs[0]
set_run(p, "Temper does not make AI smarter. It makes AI's mistakes ", 16, BODY)
set_run(p, "structurally harmless", 16, BLACK, bold=True)
set_run(p, ".", 16, BODY)

# ── SLIDE 9d: Code Review at Scale ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "AI + Code Review")
add_title(slide, "Review once, ship to every platform")
add_body(slide, "The bottleneck in AI-driven development is not code generation. It is code review. Temper compresses the review surface.", top=Inches(1.8), size=16)

add_big_num(slide, "12\u00d7", "review burden today\n(one per language)", Inches(2.0), Inches(2.6), Inches(3.5))

tf_arrow = add_textbox(slide, Inches(5.8), Inches(2.8), Inches(1.5), Inches(0.6))
add_para(tf_arrow, "\u2192", size=36, color=BORDER, align=PP_ALIGN.CENTER)

add_big_num(slide, "1\u00d7", "review burden with Temper\n(one source of truth)", Inches(7.8), Inches(2.6), Inches(3.5))

bullets = [
    "AI generates a Temper library \u2014 human reviews one codebase",
    "Approved code compiles to all targets with identical security properties",
    "No per-language security specialists needed for generated output",
    "Audit trail is one git history, not twelve",
]
tf = add_textbox(slide, MARGIN + Inches(0.2), Inches(4.6), CONTENT_W - Inches(0.2), Inches(1.5))
for b in bullets:
    add_para(tf, "\u2014  " + b, size=14, color=BODY, space_after=4)

tf2 = add_body(slide, "", top=Inches(6.0))
p = tf2.paragraphs[0]
set_run(p, "When your organization ships SDKs in 12 languages, the difference between reviewing 12 AI-generated codebases and reviewing 1 is the difference between ", 16, BODY)
set_run(p, "possible and impossible", 16, BLACK, bold=True)
set_run(p, ".", 16, BODY)

# ── SLIDE 10: AWS Scale ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "At Scale")
add_title(slide, "What if SDK security was a single source of truth?")

add_big_num(slide, "3,600", "libraries to patch today", Inches(2.0), Inches(2.2), Inches(3.5))

tf_arrow = add_textbox(slide, Inches(5.8), Inches(2.4), Inches(1.5), Inches(0.6))
add_para(tf_arrow, "\u2192", size=36, color=BORDER, align=PP_ALIGN.CENTER)

add_big_num(slide, "300", "Temper sources to maintain", Inches(7.8), Inches(2.2), Inches(3.5))

add_body(slide, "AWS manages 300+ services across 12 languages. A single vulnerability requires coordinated patches across:", top=Inches(4.0), size=16)

bullets = [
    "12 SDK repositories with independent release cycles",
    "8+ package registries (npm, PyPI, Maven, NuGet, crates.io, RubyGems, Packagist, Go proxy)",
    "Daily releases that must be synchronized",
]
tf = add_textbox(slide, MARGIN + Inches(0.2), Inches(4.7), CONTENT_W - Inches(0.2), Inches(1.5))
for b in bullets:
    add_para(tf, "\u2014  " + b, size=14, color=BODY, space_after=4)

tf2 = add_body(slide, "", top=Inches(6.0))
p = tf2.paragraphs[0]
set_run(p, "With Temper, security logic lives in ", 16, BODY)
set_run(p, "one place", 16, BLACK, bold=True)
set_run(p, ". A fix compiles to all targets in a single build. No coordination. No drift.", 16, BODY)

# ── SLIDE 11: Beyond ORMs ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "Beyond This Demo")
add_title(slide, "Where Temper changes the game")

card_data = [
    ("SDK Libraries", "API clients, serialization, auth flows, retry logic \u2014 write once, ship to every language your customers use."),
    ("Security Primitives", "Crypto wrappers, input validators, token parsers \u2014 the code that must be correct everywhere gets one audit, not twelve."),
    ("Data Validation", "Schema validators, sanitizers, format parsers \u2014 consistent behavior across frontend, backend, and systems languages."),
    ("AI-Safe Foundations", "Typed guardrails that prevent entire categories of bugs regardless of what generates the code \u2014 human or AI."),
]

card_w = Inches(5.6)
card_h = Inches(2.0)
positions = [
    (MARGIN, Inches(2.0)),
    (MARGIN + card_w + Inches(0.6), Inches(2.0)),
    (MARGIN, Inches(4.3)),
    (MARGIN + card_w + Inches(0.6), Inches(4.3)),
]

for (cx, cy), (title, desc) in zip(positions, card_data):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, card_w, card_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.75)

    tf = add_textbox(slide, cx + Inches(0.25), cy + Inches(0.2), card_w - Inches(0.5), Inches(0.4))
    add_para(tf, title, size=16, color=DARK, bold=True)

    tf2 = add_textbox(slide, cx + Inches(0.25), cy + Inches(0.7), card_w - Inches(0.5), Inches(1.2))
    add_para(tf2, desc, size=13, color=BODY)

# ── SLIDE 12: Summary ──
slide = prs.slides.add_slide(prs.slide_layouts[6])

tf = add_textbox(slide, MARGIN, Inches(0.5), CONTENT_W, Inches(0.7))
add_para(tf, "The numbers", size=36, color=DARK, bold=True, align=PP_ALIGN.CENTER)

stats_row1 = [("0", "SQL injection vulnerabilities\nacross 6 backends"), ("1", "commit to fix 3 findings\nin all 6 languages"), ("56", "tests passing\nafter remediation")]
stats_row2 = [("9", "critical vulns in conventional\nAI-generated ORM"), ("6", "language targets\nfrom single source"), ("5", "compile-time\ndefense layers")]

gap = Inches(3.8)
start = (SLIDE_W - 3 * gap) // 2

for i, (num, lbl) in enumerate(stats_row1):
    add_big_num(slide, num, lbl, start + i * gap, Inches(1.5), Inches(3.2))

for i, (num, lbl) in enumerate(stats_row2):
    add_big_num(slide, num, lbl, start + i * gap, Inches(3.5), Inches(3.2))

tf2 = add_textbox(slide, MARGIN, Inches(5.5), CONTENT_W, Inches(0.4))
add_para(tf2, "Temper does not just reduce development surface area.", size=18, color=LABEL, align=PP_ALIGN.CENTER)

tf3 = add_textbox(slide, MARGIN, Inches(6.0), CONTENT_W, Inches(0.5))
add_para(tf3, "It makes security a compile-time guarantee, not a runtime hope.", size=22, color=BLACK, bold=True, align=PP_ALIGN.CENTER)

tf4 = add_textbox(slide, MARGIN, Inches(6.8), CONTENT_W, Inches(0.4))
p = tf4.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p, "github.com/temperlang/temper", 11, LIGHT)
set_run(p, "  \u2022  ", 11, LIGHT)
set_run(p, "github.com/notactuallytreyanastasio/generic_orm", 11, LIGHT)

# ── Save ──
output_path = "/Users/robertgrayson/code/generic_temper_orm/deck.pptx"
prs.save(output_path)
print(f"Saved {output_path}")
